from __future__ import annotations

import csv
from pathlib import Path
from typing import Callable, Iterable, Iterator

from .config import SearchConfig
from .models import CancellationRequested, ExtractionResult, TextChunk


CancelCheck = Callable[[], bool]


def _split_text(text: str, location: str, chunk_chars: int, overlap: int) -> Iterator[TextChunk]:
    cleaned = text.replace("\x00", "").strip()
    if not cleaned:
        return
    if len(cleaned) <= chunk_chars:
        yield TextChunk(location=location, content=cleaned)
        return
    step = max(1, chunk_chars - min(overlap, chunk_chars // 3))
    part = 1
    for start in range(0, len(cleaned), step):
        segment = cleaned[start : start + chunk_chars].strip()
        if not segment:
            continue
        yield TextChunk(location=f"{location}・部分{part}", content=segment)
        part += 1
        if start + chunk_chars >= len(cleaned):
            break


def _limit_chunks(
    chunks: Iterable[TextChunk],
    max_chars: int,
    should_cancel: CancelCheck | None = None,
) -> tuple[list[TextChunk], bool]:
    accepted: list[TextChunk] = []
    total = 0
    truncated = False
    for chunk in chunks:
        if should_cancel and should_cancel():
            raise CancellationRequested()
        remaining = max_chars - total
        if remaining <= 0:
            truncated = True
            break
        if len(chunk.content) > remaining:
            accepted.append(TextChunk(location=chunk.location, content=chunk.content[:remaining]))
            total += remaining
            truncated = True
            break
        accepted.append(chunk)
        total += len(chunk.content)
    return accepted, truncated


def _decode_text_file(path: Path) -> tuple[str, str]:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "cp932", "utf-16"):
        try:
            return raw.decode(encoding), encoding
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace"), "utf-8-replace"


def _extract_pdf(path: Path, config: SearchConfig) -> Iterable[TextChunk]:
    import pypdfium2 as pdfium

    document = pdfium.PdfDocument(path)
    try:
        for page_number in range(1, len(document) + 1):
            page = document[page_number - 1]
            text_page = page.get_textpage()
            try:
                text = text_page.get_text_range() or ""
            finally:
                text_page.close()
                page.close()
            yield from _split_text(
                text,
                f"{page_number}ページ",
                config.chunk_chars,
                config.chunk_overlap_chars,
            )
    finally:
        document.close()


def _excel_row_chunks(
    sheet_name: str,
    rows: Iterable[tuple[int, list[object]]],
    config: SearchConfig,
    row_limit: int | None = None,
) -> Iterator[TextChunk]:
    lines: list[str] = []
    start_row = 0
    end_row = 0
    current_chars = 0

    for row_number, values in rows:
        if row_limit is not None and row_number > row_limit:
            break
        line = "\t".join("" if value is None else str(value) for value in values).strip()
        if not line:
            continue
        if lines and current_chars + len(line) + 1 > config.chunk_chars:
            yield TextChunk(
                location=f"シート「{sheet_name}」・行{start_row}～{end_row}",
                content="\n".join(lines),
            )
            lines = []
            current_chars = 0
        if not lines:
            start_row = row_number
        lines.append(line)
        end_row = row_number
        current_chars += len(line) + 1

    if lines:
        yield TextChunk(
            location=f"シート「{sheet_name}」・行{start_row}～{end_row}",
            content="\n".join(lines),
        )


def _extract_xlsx(path: Path, config: SearchConfig, row_limit: int | None = None) -> Iterable[TextChunk]:
    import openpyxl

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    try:
        for sheet in workbook.worksheets:
            rows = (
                (row_number, list(values))
                for row_number, values in enumerate(sheet.iter_rows(values_only=True), 1)
            )
            yield from _excel_row_chunks(sheet.title, rows, config, row_limit)
    finally:
        workbook.close()


def _extract_xls(path: Path, config: SearchConfig, row_limit: int | None = None) -> Iterable[TextChunk]:
    import xlrd

    workbook = xlrd.open_workbook(path, on_demand=True)
    try:
        for sheet in workbook.sheets():
            rows = ((index + 1, sheet.row_values(index)) for index in range(sheet.nrows))
            yield from _excel_row_chunks(sheet.name, rows, config, row_limit)
    finally:
        workbook.release_resources()


def _extract_docx(path: Path, config: SearchConfig) -> Iterable[TextChunk]:
    from docx import Document

    document = Document(path)
    units: list[str] = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table_number, table in enumerate(document.tables, 1):
        table_lines = []
        for row in table.rows:
            table_lines.append("\t".join(cell.text.strip() for cell in row.cells))
        if table_lines:
            units.append(f"[表{table_number}]\n" + "\n".join(table_lines))
    yield from _split_text(
        "\n\n".join(units),
        "本文・表",
        config.chunk_chars,
        config.chunk_overlap_chars,
    )


def _extract_pptx(path: Path, config: SearchConfig) -> Iterable[TextChunk]:
    from pptx import Presentation

    presentation = Presentation(path)
    for slide_number, slide in enumerate(presentation.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                texts.append(text.strip())
        yield from _split_text(
            "\n".join(texts),
            f"スライド{slide_number}",
            config.chunk_chars,
            config.chunk_overlap_chars,
        )


def _extract_delimited(path: Path, config: SearchConfig, row_limit: int | None = None) -> Iterable[TextChunk]:
    text, encoding = _decode_text_file(path)
    lines = text.splitlines()
    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
    reader = csv.reader(lines, delimiter=delimiter)
    rows = ((number, row) for number, row in enumerate(reader, 1))
    yield from _excel_row_chunks(f"CSV ({encoding})", rows, config, row_limit)


def _extract_plain_text(path: Path, config: SearchConfig) -> Iterable[TextChunk]:
    text, encoding = _decode_text_file(path)
    yield from _split_text(
        text,
        f"テキスト ({encoding})",
        config.chunk_chars,
        config.chunk_overlap_chars,
    )


def extract_file(
    path: Path,
    config: SearchConfig,
    should_cancel: CancelCheck | None = None,
) -> ExtractionResult:
    try:
        if should_cancel and should_cancel():
            raise CancellationRequested()
        size = path.stat().st_size
        policy = config.file_policy(path)
        if policy == "metadata":
            return ExtractionResult(status="metadata_only", error="設定によりファイル名のみ索引しました")
        if size > config.max_file_size_bytes and not (
            config.has_file_policy(path) and policy == "full"
        ):
            return ExtractionResult(
                status="too_large",
                error=f"ファイルサイズが上限を超えています: {size} bytes",
            )

        extension = path.suffix.lower()
        extractors = {
            ".pdf": _extract_pdf,
            ".xlsx": _extract_xlsx,
            ".xls": _extract_xls,
            ".docx": _extract_docx,
            ".pptx": _extract_pptx,
            ".csv": _extract_delimited,
            ".txt": _extract_plain_text,
        }
        extractor = extractors.get(extension)
        if extractor is None:
            return ExtractionResult(status="unsupported", error=f"未対応形式: {extension}")

        extracted = (
            extractor(path, config, config.table_head_rows)
            if policy == "head" and extension in {".xlsx", ".xls", ".csv"}
            else extractor(path, config)
        )
        chunks, truncated = _limit_chunks(
            extracted,
            config.max_text_chars_per_file,
            should_cancel,
        )
        if not chunks:
            return ExtractionResult(status="empty", error="本文を抽出できませんでした")
        head_limited = policy == "head" and extension in {".xlsx", ".xls", ".csv"}
        return ExtractionResult(
            chunks=chunks,
            status="truncated" if truncated or head_limited else "ok",
            error=(
                f"設定により先頭{config.table_head_rows:,}行だけ索引しました"
                if head_limited
                else "本文が設定上限に達したため切り詰めました" if truncated else ""
            ),
            truncated=truncated or head_limited,
        )
    except CancellationRequested:
        raise
    except Exception as exc:  # extraction errors must not stop the complete index run
        return ExtractionResult(status="error", error=f"{type(exc).__name__}: {exc}")
